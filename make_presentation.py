"""Generate the Two-Clocks lecture presentation (PPTX).

Redesigned for modern lecture delivery:
  * Visual clarity / low cognitive load  -> one idea per slide, big bold
    geometric type, generous white space, soft rounded geometry.
  * Restorative palette                  -> warm neutrals + calm teal + earthy
    olive, warm clay instead of alarm-red.
  * Non-linear delivery                  -> a hyperlinked "menu" slide and a
    persistent Menu pill on every slide (click straight to any section).
  * Progressive disclosure               -> key concepts are built up across
    staged slides (slow clock -> fast clock; GLE -> mutilation -> repair).
  * Data as narrative                    -> every chart slide has an active
    takeaway title + one highlighted number (one chart = one insight).
  * Audience integration                 -> Slido/Mentimeter poll slides and
    "spacing effect" think-pair-share review slides at section breaks.
  * Trigger framing                      -> "watch & decide" dilemma slides that
    pose the problem before the theory (speaker notes carry the clip cue).

Every scientific claim and number is taken from the repo (README + REPORT_*.md);
nothing in the original deck's content is dropped.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_CONNECTOR
from PIL import Image
import os
import warnings

# ---------------------------------------------------------------------------
# Restorative palette (warm neutrals, calm teal, earthy olive, warm clay)
# ---------------------------------------------------------------------------
CREAM      = RGBColor(0xF5, 0xF0, 0xE6)   # warm neutral page
CARD       = RGBColor(0xFC, 0xFA, 0xF4)   # lighter card / figure mat
SAND       = RGBColor(0xEA, 0xE1, 0xCE)   # soft panel
INK        = RGBColor(0x33, 0x2E, 0x29)   # warm near-black text
MUTE       = RGBColor(0x8A, 0x80, 0x72)   # muted secondary text
TEAL       = RGBColor(0x2A, 0x8C, 0x82)   # calm teal (primary accent)
TEAL_DEEP  = RGBColor(0x1E, 0x5F, 0x59)   # deep teal (section bands)
OLIVE      = RGBColor(0x7C, 0x8A, 0x52)   # earthy green (success / FDT)
CLAY       = RGBColor(0xC0, 0x6B, 0x4F)   # warm terracotta (failure, used softly)
OCHRE      = RGBColor(0xD2, 0xA1, 0x4A)   # muted gold (audience / interaction)
CREAM_SOFT = RGBColor(0xEC, 0xE9, 0xDF)   # divider text-muted on dark

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN  = 0.9
CW      = 13.333 - 2 * MARGIN  # content width

HEAD = "Century Gothic"   # geometric sans for display / titles / numbers
BODY = "Segoe UI"         # clean sans for body + captions
MONO = "Consolas"         # equations

_page = {"n": 0}


# ---------------------------------------------------------------------------
# Low-level helpers
# ---------------------------------------------------------------------------
def _blank(prs, bg=CREAM):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = bg
    return slide


def _no_line(shape):
    shape.line.fill.background()


def _no_shadow(shape):
    shape.shadow.inherit = False


def rrect(slide, l, t, w, h, fill, radius=0.10, line=None, line_w=1.0):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(l), Inches(t), Inches(w), Inches(h))
    try:
        sp.adjustments[0] = radius
    except (IndexError, TypeError):
        pass
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        _no_line(sp)
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    _no_shadow(sp)
    return sp


def circle(slide, l, t, d, fill, line=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                Inches(l), Inches(t), Inches(d), Inches(d))
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        _no_line(sp)
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1.0)
    _no_shadow(sp)
    return sp


def hrule(slide, l, t, w, color=TEAL, weight=2.5):
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                    Inches(l), Inches(t), Inches(l + w), Inches(t))
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    _no_shadow(ln)
    return ln


def _set_run(r, text, size, color, bold, font, italic=False):
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font


def textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tf


def line_para(tf, runs, size=20, color=INK, bold=False, font=BODY,
              align=PP_ALIGN.LEFT, space_before=0.0, space_after=0.0,
              line_spacing=1.0, first=False):
    """runs: str OR list of (text, color, bold) tuples."""
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    if space_after:
        p.space_after = Pt(space_after)
    p.line_spacing = line_spacing
    if isinstance(runs, str):
        runs = [(runs, color, bold)]
    for tup in runs:
        txt, c, b = (tup + (bold,))[:3] if len(tup) == 2 else tup
        r = p.add_run()
        _set_run(r, txt, size, c, b, font)
    return p


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def footer(slide, menu_slide=None, light=False):
    """Page number + repo handle + (optional) Menu pill for non-linear nav."""
    _page["n"] += 1
    col = CREAM_SOFT if light else MUTE
    tf = textbox(slide, MARGIN, 7.06, 6.0, 0.3)
    line_para(tf, [("Abdelrahman Saifelden", col, False)], size=10, font=BODY, first=True)
    tf2 = textbox(slide, 9.0, 7.06, 3.43 - 0.0, 0.3)
    line_para(tf2, [(f"{_page['n']:02d}", col, True)], size=10, font=BODY,
              align=PP_ALIGN.RIGHT, first=True)
    if menu_slide is not None:
        pill = rrect(slide, 11.45, 7.0, 0.98, 0.34,
                     SAND if not light else TEAL_DEEP, radius=0.5)
        ptf = pill.text_frame
        ptf.word_wrap = False
        ptf.vertical_anchor = MSO_ANCHOR.MIDDLE
        ptf.margin_left = 0
        ptf.margin_right = 0
        line_para(ptf, [("\u2630  Menu", INK if not light else CREAM, True)],
                  size=11, font=BODY, align=PP_ALIGN.CENTER, first=True)
        pill.click_action.target_slide = menu_slide


def kicker(slide, text, color=TEAL, top=0.62, left=MARGIN):
    tf = textbox(slide, left, top, CW, 0.32)
    line_para(tf, [(text.upper(), color, True)], size=13, font=BODY, first=True)


def place_figure(slide, path, l, t, w, h, mat=True, pad=0.22):
    """Aspect-fit a figure inside a soft white mat, centered in (l,t,w,h)."""
    path = _FigDirs.resolve(path)
    if mat:
        rrect(slide, l, t, w, h, CARD, radius=0.045)
    try:
        with Image.open(path) as im:
            pw, ph = im.size
    except FileNotFoundError:
        warnings.warn(f"Figure not found: {path!r}; skipping (mat placeholder kept)")
        return
    except OSError:
        warnings.warn(f"Cannot read figure {path!r}; skipping (mat placeholder kept)")
        return
    aspect = pw / ph
    aw, ah = w - 2 * pad, h - 2 * pad
    if aw / ah > aspect:
        fh = ah
        fw = fh * aspect
    else:
        fw = aw
        fh = fw / aspect
    fl = l + (w - fw) / 2
    ft = t + (h - fh) / 2
    slide.shapes.add_picture(path, Inches(fl), Inches(ft), Inches(fw), Inches(fh))


def stat_pill(slide, l, t, w, label, value, accent=TEAL, h=1.0):
    """Soft card highlighting one number (data-as-narrative)."""
    rrect(slide, l, t, w, h, CARD, radius=0.16)
    rrect(slide, l, t, 0.12, h, accent, radius=0.16)
    tf = textbox(slide, l + 0.30, t + 0.12, w - 0.45, h - 0.2, MSO_ANCHOR.MIDDLE)
    line_para(tf, [(value, accent, True)], size=27, font=HEAD, first=True)
    line_para(tf, [(label, MUTE, False)], size=12.5, font=BODY, space_before=2)


# ---------------------------------------------------------------------------
# Slide templates
# ---------------------------------------------------------------------------
def title_slide(prs):
    slide = _blank(prs, CREAM)
    # soft geometric accents
    circle(slide, -1.6, -1.6, 4.4, SAND)
    circle(slide, 11.0, 5.0, 4.6, RGBColor(0xEE, 0xE7, 0xD6))
    rrect(slide, MARGIN, 1.5, 0.9, 0.16, TEAL, radius=0.5)
    tf = textbox(slide, MARGIN, 1.85, 11.2, 3.2)
    line_para(tf, [("Two Clocks", TEAL, True)], size=64, font=HEAD, bold=True, first=True)
    line_para(tf, [("Diagnosing & repairing K-theory's", INK, True)], size=33,
              font=HEAD, space_before=8)
    line_para(tf, [("structural blindness", INK, True)], size=33, font=HEAD)
    sub = textbox(slide, MARGIN, 5.15, 11.0, 1.4)
    line_para(sub, [("Pressure and temperature run on two different clocks. ",
                     INK, False),
                    ("One closure that ignores the fast one mutilates "
                     "Navier\u2013Stokes \u2014 and there is an exact repair.",
                     MUTE, False)],
              size=18, font=BODY, line_spacing=1.15, first=True)
    auth = textbox(slide, MARGIN, 6.45, 11.5, 0.9)
    line_para(auth, [("Abdelrahman Saifelden", INK, True)], size=20, font=HEAD,
              bold=True, first=True)
    line_para(auth, [("Mechatronics Engineering \u00b7 AIET, Alexandria",
                      MUTE, False)], size=13, font=BODY, space_before=2)
    notes(slide, "Open cold. Do NOT read the slide. One line: 'For a century we "
                 "stirred pressure and temperature with one number. Tonight: why "
                 "that quietly breaks the equations \u2014 and the exact fix.' "
                 "Then go to the trigger clip.")
    return slide


def trigger_slide(prs, menu, kick, headline_runs, prompt, clip_cue, accent=OCHRE):
    """'Watch & decide' dilemma slide \u2014 pose the problem before the theory."""
    slide = _blank(prs, CREAM)
    kicker(slide, kick, color=accent)
    # play badge
    circle(slide, MARGIN, 1.25, 0.95, accent)
    tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                                 Inches(MARGIN + 0.34), Inches(1.5),
                                 Inches(0.34), Inches(0.45))
    tri.rotation = 90
    tri.fill.solid(); tri.fill.fore_color.rgb = CREAM; _no_line(tri); _no_shadow(tri)
    tf = textbox(slide, MARGIN + 1.25, 1.2, CW - 1.25, 1.7, MSO_ANCHOR.MIDDLE)
    for ln in headline_runs:
        line_para(tf, ln, size=34, font=HEAD, bold=True, line_spacing=1.04,
                  first=(ln is headline_runs[0]))
    # the dilemma
    panel = rrect(slide, MARGIN, 3.55, CW, 2.5, CARD, radius=0.06)
    ptf = textbox(slide, MARGIN + 0.5, 3.85, CW - 1.0, 1.9, MSO_ANCHOR.MIDDLE)
    line_para(ptf, [("Before the theory \u2014 think:", accent, True)], size=15,
              font=BODY, first=True)
    line_para(ptf, [(prompt, INK, False)], size=22, font=BODY, space_before=8,
              line_spacing=1.12)
    foot = textbox(slide, MARGIN, 6.25, CW, 0.5)
    line_para(foot, [("\u25b6  ", accent, True), (clip_cue, MUTE, False)],
              size=13, font=BODY, first=True)
    footer(slide, menu)
    notes(slide, "TRIGGER VIDEO: " + clip_cue + " Play the 30\u201390 s clip, then "
                 "pause. Ask the room the question on screen. Take 2\u20133 verbal "
                 "guesses (or run the poll on the next slide) BEFORE revealing any "
                 "theory. The goal is to make them want the answer.")
    return slide


def poll_slide(prs, menu, question, options, code="TWOCLOCKS", tool="Slido"):
    """Audience Response System slide (Mentimeter / Slido / ClassPoint)."""
    slide = _blank(prs, CREAM)
    rrect(slide, 0, 0, 13.333, 1.5, TEAL_DEEP, radius=0.0)
    kt = textbox(slide, MARGIN, 0.42, CW, 0.4)
    line_para(kt, [("LIVE POLL  \u00b7  PHONES OUT", OCHRE, True)], size=14,
              font=BODY, first=True)
    jt = textbox(slide, MARGIN, 0.78, CW, 0.6)
    line_para(jt, [(f"Join at {tool.lower()}.com", CREAM, True),
                   ("    \u00b7    code ", CREAM_SOFT, False),
                   (f"#{code}", OCHRE, True)], size=20, font=HEAD, first=True)
    qt = textbox(slide, MARGIN, 1.95, CW, 1.4, MSO_ANCHOR.TOP)
    line_para(qt, [(question, INK, True)], size=30, font=HEAD, bold=True,
              line_spacing=1.05, first=True)
    # option cards
    n = len(options)
    top = 3.5
    gap = 0.25
    ch = (6.7 - top - (n - 1) * gap) / n
    letters = "ABCDEF"
    for i, opt in enumerate(options):
        ot = top + i * (ch + gap)
        rrect(slide, MARGIN, ot, CW, ch, CARD, radius=0.18)
        circle(slide, MARGIN + 0.22, ot + (ch - 0.62) / 2, 0.62, SAND)
        lt = textbox(slide, MARGIN + 0.22, ot + (ch - 0.62) / 2, 0.62, 0.62,
                     MSO_ANCHOR.MIDDLE)
        line_para(lt, [(letters[i], TEAL, True)], size=20, font=HEAD,
                  align=PP_ALIGN.CENTER, first=True)
        ttf = textbox(slide, MARGIN + 1.1, ot, CW - 1.4, ch, MSO_ANCHOR.MIDDLE)
        line_para(ttf, [(opt, INK, False)], size=19, font=BODY, first=True)
    footer(slide, menu)
    notes(slide, "AUDIENCE RESPONSE (Slido/Mentimeter/ClassPoint). Let votes land "
                 "live; do not editorialise yet. Whatever wins, the next slides "
                 "adjudicate it with computed evidence. Forces active recall + "
                 "gives you instant read on comprehension.")
    return slide


def review_slide(prs, menu, headline, task, accent=OLIVE):
    """Spacing-effect micro-review: think-pair-share, ~60 s."""
    slide = _blank(prs, CREAM)
    rrect(slide, MARGIN, 1.4, CW, 4.7, CARD, radius=0.05)
    rrect(slide, MARGIN, 1.4, CW, 0.86, accent, radius=0.07)
    kt = textbox(slide, MARGIN + 0.5, 1.5, CW - 1.0, 0.66, MSO_ANCHOR.MIDDLE)
    line_para(kt, [("\u23f1  60-SECOND REVIEW  \u00b7  THINK \u2013 PAIR \u2013 SHARE",
                    CREAM, True)], size=15, font=BODY, first=True)
    ht = textbox(slide, MARGIN + 0.5, 2.7, CW - 1.0, 1.4)
    line_para(ht, [(headline, INK, True)], size=28, font=HEAD, bold=True,
              line_spacing=1.05, first=True)
    tt = textbox(slide, MARGIN + 0.5, 4.4, CW - 1.0, 1.4)
    line_para(tt, [("Your turn:  ", accent, True), (task, INK, False)],
              size=20, font=BODY, line_spacing=1.15, first=True)
    footer(slide, menu)
    notes(slide, "SPACING EFFECT. Stop the monologue. 20 s think, 30 s pair, "
                 "10 s share-out. Distributed retrieval beats a 45-min lecture "
                 "with one summary at the end. Don't move on until two pairs "
                 "answer out loud.")
    return slide


def section_divider(prs, menu, num, title_runs, subtitle):
    slide = _blank(prs, TEAL_DEEP)
    circle(slide, 9.6, -1.9, 5.4, RGBColor(0x24, 0x6B, 0x64))
    circle(slide, 10.8, 4.4, 4.0, RGBColor(0x20, 0x66, 0x5F))
    nt = textbox(slide, MARGIN, 1.9, 4.0, 2.0)
    line_para(nt, [(num, OCHRE, True)], size=120, font=HEAD, bold=True, first=True)
    hrule(slide, MARGIN, 4.05, 3.2, color=OCHRE, weight=3)
    tt = textbox(slide, MARGIN, 4.25, 11.0, 1.7)
    for i, ln in enumerate(title_runs):
        line_para(tt, ln, size=40, font=HEAD, bold=True, line_spacing=1.02,
                  first=(i == 0))
    st = textbox(slide, MARGIN, 6.05, 10.5, 0.8)
    line_para(st, [(subtitle, CREAM_SOFT, False)], size=17, font=BODY,
              line_spacing=1.1, first=True)
    footer(slide, menu, light=True)
    return slide


def statement_slide(prs, menu, kick, lines, sub=None, accent=TEAL):
    """One big idea. Minimal text, maximal white space."""
    slide = _blank(prs, CREAM)
    kicker(slide, kick, color=accent)
    tf = textbox(slide, MARGIN, 2.3, CW, 3.0, MSO_ANCHOR.TOP)
    for i, ln in enumerate(lines):
        line_para(tf, ln, size=46, font=HEAD, bold=True, line_spacing=1.04,
                  first=(i == 0))
    if sub:
        stf = textbox(slide, MARGIN, 5.35, CW - 1.0, 1.4)
        line_para(stf, [(sub, MUTE, False)], size=19, font=BODY,
                  line_spacing=1.18, first=True)
    footer(slide, menu)
    return slide


def figure_slide(prs, menu, kick, takeaway_runs, fig_path, stats=None,
                 caption=None, accent=TEAL):
    """One chart = one insight. Active takeaway title + highlighted number."""
    slide = _blank(prs, CREAM)
    kicker(slide, kick, color=accent)
    tt = textbox(slide, MARGIN, 0.95, CW, 1.25)
    for i, ln in enumerate(takeaway_runs):
        line_para(tt, ln, size=27, font=HEAD, bold=True, line_spacing=1.04,
                  first=(i == 0))
    has_stats = bool(stats)
    fig_w = (CW - 3.4) if has_stats else CW
    fig_h = 4.15
    place_figure(slide, fig_path, MARGIN, 2.3, fig_w, fig_h)
    if has_stats:
        sx = MARGIN + fig_w + 0.3
        sw = 13.333 - MARGIN - sx
        n = len(stats)
        sh = 1.0
        gap = (fig_h - n * sh) / max(n - 1, 1) if n > 1 else 0
        for i, (label, value, acc) in enumerate(stats):
            stat_pill(slide, sx, 2.3 + i * (sh + gap), sw, label, value, acc, h=sh)
    if caption:
        ct = textbox(slide, MARGIN, 6.58, CW - 0.2, 0.4)
        line_para(ct, [(caption, MUTE, False)], size=12, font=BODY,
                  line_spacing=1.0, first=True)
    footer(slide, menu)
    return slide


def cards_slide(prs, menu, kick, title_runs, cards, accent=TEAL, cols=2):
    """Soft-geometry cards \u2014 low-clutter alternative to a bullet wall."""
    slide = _blank(prs, CREAM)
    kicker(slide, kick, color=accent)
    tt = textbox(slide, MARGIN, 0.95, CW, 1.1)
    for i, ln in enumerate(title_runs):
        line_para(tt, ln, size=27, font=HEAD, bold=True, line_spacing=1.03,
                  first=(i == 0))
    n = len(cards)
    rows = (n + cols - 1) // cols
    gx, gy = 0.3, 0.3
    top = 2.25
    cw = (CW - (cols - 1) * gx) / cols
    ch = (6.75 - top - (rows - 1) * gy) / rows
    for i, (tag, body, acc) in enumerate(cards):
        r, c = divmod(i, cols)
        cl = MARGIN + c * (cw + gx)
        ct = top + r * (ch + gy)
        rrect(slide, cl, ct, cw, ch, CARD, radius=0.07)
        rrect(slide, cl, ct, cw, 0.10, acc, radius=0.07)
        tf = textbox(slide, cl + 0.32, ct + 0.26, cw - 0.6, ch - 0.5)
        line_para(tf, [(tag, acc, True)], size=17, font=HEAD, bold=True, first=True)
        line_para(tf, [(body, INK, False)], size=14.5, font=BODY, space_before=7,
                  line_spacing=1.12)
    footer(slide, menu)
    return slide


def equation_slide(prs, menu, kick, title_runs, eq_runs_lines, notes_lines,
                   accent=TEAL):
    """Centered equation card + short, plain-language gloss (built up if staged)."""
    slide = _blank(prs, CREAM)
    kicker(slide, kick, color=accent)
    tt = textbox(slide, MARGIN, 0.95, CW, 1.0)
    for i, ln in enumerate(title_runs):
        line_para(tt, ln, size=26, font=HEAD, bold=True, line_spacing=1.03,
                  first=(i == 0))
    rrect(slide, MARGIN, 2.15, CW, 1.95, INK if False else TEAL_DEEP, radius=0.07)
    etf = textbox(slide, MARGIN + 0.5, 2.25, CW - 1.0, 1.75, MSO_ANCHOR.MIDDLE)
    for i, ln in enumerate(eq_runs_lines):
        line_para(etf, ln, size=22, font=MONO, align=PP_ALIGN.CENTER,
                  line_spacing=1.15, first=(i == 0))
    ntf = textbox(slide, MARGIN, 4.45, CW, 2.3)
    for i, ln in enumerate(notes_lines):
        line_para(ntf, ln, size=17, font=BODY, space_before=(0 if i == 0 else 7),
                  line_spacing=1.12, first=(i == 0))
    footer(slide, menu)
    return slide


def menu_slide(prs):
    """Hyperlinked home menu \u2014 non-linear navigation. Targets wired later."""
    slide = _blank(prs, CREAM)
    kicker(slide, "Choose your path \u00b7 click any card", color=TEAL)
    tt = textbox(slide, MARGIN, 0.95, CW, 1.0)
    line_para(tt, [("Where should we start?", INK, True)], size=34, font=HEAD,
              bold=True, first=True)
    st = textbox(slide, MARGIN, 1.7, CW, 0.5)
    line_para(st, [("This deck is non-linear \u2014 jump to whatever the room "
                    "wants first. The ", MUTE, False),
                   ("\u2630 Menu", TEAL, True),
                   (" pill on every slide brings you back here.", MUTE, False)],
              size=14, font=BODY, first=True)
    cards = [
        ("01", "The two clocks", "The one idea everything rests on", TEAL),
        ("02", "The diagnosis", "Real data: where K-theory cracks", CLAY),
        ("03", "The evidence", "Projection vs. spectrum-matching", OLIVE),
        ("04", "The repair", "Mori\u2013Zwanzig / projected-FDT", TEAL_DEEP),
        ("05", "The benchmark", "The test that decides it", OCHRE),
        ("06", "Why it matters", "Four real-world fast-clock traps", CLAY),
    ]
    cols = 3
    gx, gy = 0.35, 0.35
    top = 2.55
    cw = (CW - (cols - 1) * gx) / cols
    ch = (6.7 - top - gy) / 2
    handles = []
    for i, (num, title, body, acc) in enumerate(cards):
        r, c = divmod(i, cols)
        cl = MARGIN + c * (cw + gx)
        ct = top + r * (ch + gy)
        card = rrect(slide, cl, ct, cw, ch, CARD, radius=0.08)
        rrect(slide, cl, ct, 0.12, ch, acc, radius=0.08)
        tf = textbox(slide, cl + 0.4, ct + 0.28, cw - 0.7, ch - 0.5)
        line_para(tf, [(num, acc, True)], size=26, font=HEAD, bold=True, first=True)
        line_para(tf, [(title, INK, True)], size=19, font=HEAD, bold=True,
                  space_before=4)
        line_para(tf, [(body, MUTE, False)], size=13, font=BODY, space_before=5,
                  line_spacing=1.1)
        handles.append(card)
    footer(slide, None)
    notes(slide, "NON-LINEAR DELIVERY. Ask: 'Which one should we tackle first?' "
                 "and click straight to it. Makes the lecture feel like a tailored "
                 "conversation. You can always return via the Menu pill.")
    return slide, handles


# ---------------------------------------------------------------------------
# Build the deck
# ---------------------------------------------------------------------------

class _FigDirs:
    """Resolve figure files across the reorganized domain folders.

    Keeps the legacy f"{fig}/<name>.png" call sites working: formatting the
    object yields the sentinel "<FIGS>/<name>.png", which place_figure()
    resolves against each domain's figures/ directory via resolve().
    """

    _DIRS = ("general_two_clocks/figures", "atmosphere/figures",
             "glaciers/figures", "ocean/figures", "figures")

    def __str__(self):
        return "<FIGS>"

    @classmethod
    def resolve(cls, path):
        if not isinstance(path, str) or "<FIGS>" not in path:
            return path
        name = path.split("<FIGS>/", 1)[1]
        here = os.path.dirname(os.path.abspath(__file__))
        for d in cls._DIRS:
            cand = os.path.join(here, d, name)
            if os.path.exists(cand):
                return cand
        return os.path.join(here, "figures", name)


def main():
    _page["n"] = 0
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    fig = _FigDirs()
    anchors = {}

    # ---- Opening: title -> trigger -> poll -> menu -------------------------
    title_slide(prs)

    menu, menu_cards = menu_slide(prs)  # created early so footers can link to it

    trigger_slide(
        prs, menu, "The hook \u00b7 watch & decide",
        [[("In 24 hours, Hurricane Otis went", INK, True)],
         [("tropical storm ", INK, True), ("\u2192", CLAY, True),
          (" Category 5.", INK, True)]],
        "Every operational model had the data \u2014 ocean heat, steering winds, "
        "all correct \u2014 and almost none predicted the jump. What were they "
        "missing?",
        "Clip cue: 30\u201360 s of Otis rapid-intensification satellite loop (2023).",
        accent=CLAY)

    poll_slide(prs, menu,
               "Why do turbulence models miss events like this?",
               ["Not enough data / resolution",
                "Bad tuning of the coefficients",
                "Something structural \u2014 the math itself is blind",
                "Just bad luck / chaos"],
               code="TWOCLOCKS", tool="Slido")

    # ===== SECTION 01 \u2014 THE TWO CLOCKS =================================
    s = statement_slide(
        prs, menu, "01 \u00b7 The two clocks",
        [[("Pressure and temperature", INK, True)],
         [("are ", INK, True), ("not", CLAY, True), (" one coupled scalar.", INK, True)]],
        sub="They run on fundamentally different clocks \u2014 and on different "
            "spatial architectures. Stir both with a single eddy diffusivity and "
            "you erase one of them.")
    anchors["01"] = s

    statement_slide(
        prs, menu, "The slow clock",
        [[("Temperature is ", INK, True), ("local", TEAL, True), (".", INK, True)]],
        sub="A parabolic field. It diffuses neighbour-to-neighbour (the heat "
            "equation) and is torn into sharp filaments by local shear. "
            "Memoryless, gradient-driven \u2014 exactly what K-theory sees.",
        accent=TEAL)

    statement_slide(
        prs, menu, "The fast clock",
        [[("Pressure is ", INK, True), ("global", CLAY, True), (".", INK, True)]],
        sub="An elliptic field. It solves a Poisson equation \u2207\u00b2p = f that "
            "feels the whole domain and its boundaries at once \u2014 instantaneous, "
            "non-local, boundary-aware. Invisible to any local diffusion.",
        accent=CLAY)

    statement_slide(
        prs, menu, "The thesis in one line",
        [[("K-theory models ", INK, True), ("only", CLAY, True),
          (" the slow clock.", INK, True)]],
        sub="Classical single-Prandtl closure (Smagorinsky) is structurally "
            "unable to represent the fast, elliptic pressure clock. The rest of "
            "this lecture proves it \u2014 then repairs it.")

    review_slide(prs, menu,
                 "Say it back: what are the two clocks?",
                 "In one sentence to your neighbour \u2014 which field is local & "
                 "parabolic, which is global & elliptic, and which one does "
                 "K-theory ignore?")

    # ===== SECTION 02 \u2014 THE DIAGNOSIS ==================================
    d = section_divider(
        prs, menu, "02",
        [[("The diagnosis", CREAM, True)]],
        "Three independent public datasets, increasing complexity \u2014 every "
        "headline number computed by the code in the repo, not asserted.")
    anchors["02"] = d

    figure_slide(
        prs, menu, "Part 1 \u00b7 NEON eddy-flux, Wind River forest",
        [[("Temperature ticks at 24 h. Pressure carries a 12 h tide.", INK, True)]],
        f"{fig}/03_spectra_two_clocks.png",
        stats=[("Pressure\u2013temperature coupling", "r \u2248 0.07", CLAY),
               ("Variance explained", "R\u00b2 \u2248 0.005", CLAY)],
        caption="Two qualitatively different spectral fingerprints from the same "
                "site \u2014 not one coupled scalar.")

    figure_slide(
        prs, menu, "Part 1b \u00b7 ASOS 1-min, three stations 25\u201342\u00b0N",
        [[("The pressure tide weakens from equator to pole \u2014", INK, True)],
         [("it ", INK, True), ("knows about the whole planet", CLAY, True),
          (".", INK, True)]],
        f"{fig}/11_tide_vs_latitude.png",
        stats=[("S\u2082 tide  MIA (25.8\u00b0N)", "1.08 hPa", TEAL),
               ("S\u2082 tide  DFW (32.9\u00b0N)", "0.96 hPa", TEAL),
               ("S\u2082 tide  DSM (41.5\u00b0N)", "0.63 hPa", CLAY)],
        caption="Temperature peaks diurnally at every latitude; the 12 h "
                "pressure tide is a global resonance.")

    figure_slide(
        prs, menu, "Part 3 \u00b7 NEON stability analysis, 602 periods",
        [[("Transport efficiency swings ", INK, True), ("4.3\u00d7", CLAY, True),
          (" \u2014 one fixed", INK, True)],
         [("Prandtl number cannot track that.", INK, True)]],
        f"{fig}/07_transport_decoupling.png",
        stats=[("Momentum/heat ratio, strongly stable", "0.77", TEAL),
               ("\u2026 near-neutral", "3.28", CLAY),
               ("Single-K assumes this is", "constant", MUTE)],
        caption="Momentum and heat efficiencies move oppositely \u2014 a constant "
                "Prandtl number cannot represent that.")

    review_slide(prs, menu,
                 "Why can't one diffusivity capture this?",
                 "30 seconds with a neighbour \u2014 name the single assumption in "
                 "K-theory that the 4.3\u00d7 swing breaks.",
                 accent=CLAY)

    # ===== SECTION 03 \u2014 THE EVIDENCE ====================================
    e = section_divider(
        prs, menu, "03",
        [[("The evidence", CREAM, True)]],
        "Take the two clocks apart in a controlled flow, then confirm the split "
        "in real reanalysis winds. The decisive structural test lives here.")
    anchors["03"] = e

    figure_slide(
        prs, menu, "Part 6 \u00b7 Boussinesq projection step",
        [[("One Poisson solve drops divergence from ", INK, True),
          ("10\u207b\u00b2 to 10\u207b\u00b9\u2075", TEAL, True), (".", INK, True)]],
        f"{fig}/23_projection_step.png",
        stats=[("Slow drift alone, RMS \u2207\u00b7u*", "2.7\u00d710\u207b\u00b2", CLAY),
               ("After one Leray projection", "2.3\u00d710\u207b\u00b9\u2075", OLIVE),
               ("That projection is", "the fast clock", TEAL)],
        caption="The local drift breaks mass conservation; one global elliptic "
                "solve fixes it instantly \u2014 the fast clock.")

    figure_slide(
        prs, menu, "Part 6 \u00b7 the SPDE limit (the crux)",
        [[("Same spectrum, ", INK, True), ("wrong structure", CLAY, True),
          (": correlation ", INK, True), ("0.01", CLAY, True), (".", INK, True)]],
        f"{fig}/24_spde_limit.png",
        stats=[("Surrogate energy spectrum E(k)", "identical", OLIVE),
               ("Pointwise corr. with truth", "\u2248 0.012", CLAY),
               ("RMS \u2207\u00b7u (should be 0)", "3.8\u00d710\u207b\u00b2", CLAY)],
        caption="Parseval-exact spectrum, yet geometrically unrelated to the "
                "truth \u2014 energy matched, structure lost.")

    figure_slide(
        prs, menu, "Part 7 \u00b7 NCEP/NCAR reanalysis winds",
        [[("In real winds, ", INK, True), ("96%", TEAL, True),
          (" of the energy is the slow clock.", INK, True)]],
        f"{fig}/25_reanalysis_two_clocks.png",
        stats=[("KE_div / KE_rot  @ 850 hPa", "4.0%", CLAY),
               ("@ 500 hPa (non-divergence)", "0.7%", TEAL),
               ("6-hourly > daily mean", "1.1% vs 0.7%", OCHRE)],
        caption="Rotational wind is the weather; the divergent fast clock is a "
                "weak, structured residual that averaging removes.")

    poll_slide(prs, menu,
               "Predict: will the spectrum-matched surrogate pass the benchmark?",
               ["Yes \u2014 right energy means right physics",
                "No \u2014 it has the wrong structure",
                "Partly \u2014 it'll pass some tests, fail others"],
               code="TWOCLOCKS", tool="Mentimeter")

    # ===== SECTION 04 \u2014 THE REPAIR ======================================
    r = section_divider(
        prs, menu, "04",
        [[("The repair", CREAM, True)]],
        "From diagnosis to prescription: the exact resolved-scale dynamics, what "
        "K-theory deletes from it, and the closure that puts it back.")
    anchors["04"] = r

    # progressive build: GLE revealed term by term
    equation_slide(
        prs, menu, "Part 8 \u00b7 the exact dynamics (build 1 of 3)",
        [[("Start from what is ", INK, True), ("exact", TEAL, True),
          (" \u2014 no modelling yet", INK, True)]],
        [[("\u2202\u209c\u00e2 = P\u2112\u00e2", TEAL, True),
          ("  \u2212  \u222b\u2080\u1d57 K(t\u2212s) \u00e2(s) ds  +  f(t)", MUTE, False)]],
        [[("The Mori\u2013Zwanzig generalized Langevin equation. The first term is "
           "the ", INK, False), ("slow Markov drift", TEAL, True),
          (" \u2014 local, memoryless. This is the part K-theory keeps.", INK, False)]])

    equation_slide(
        prs, menu, "Part 8 \u00b7 the exact dynamics (build 2 of 3)",
        [[("The fast clock lives in the ", INK, True),
          ("memory kernel", CLAY, True)]],
        [[("\u2202\u209c\u00e2 = P\u2112\u00e2  \u2212  ", MUTE, False),
          ("\u222b\u2080\u1d57 K(t\u2212s) \u00e2(s) ds", CLAY, True),
          ("  +  f(t)", MUTE, False)]],
        [[("The convolution is ", INK, False), ("non-local in time", CLAY, True),
          (" \u2014 the resolved scales remember their past interaction with the "
           "unresolved fast modes. Delete it and the fast clock vanishes.", INK, False)]])

    equation_slide(
        prs, menu, "Part 8 \u00b7 the exact dynamics (build 3 of 3)",
        [[("\u2026 and the FDT-locked ", INK, True), ("noise", OLIVE, True)]],
        [[("\u2202\u209c\u00e2 = P\u2112\u00e2  \u2212  \u222b\u2080\u1d57 K(t\u2212s)\u00e2(s)ds  +  ",
           MUTE, False), ("f(t)", OLIVE, True)],
         [("K(\u03c4) = \u27e8f(\u03c4) f(0)\u27e9 \u00b7 \u27e8\u00e2,\u00e2\u27e9\u207b\u00b9",
           OLIVE, True)]],
        [[("The 2nd fluctuation\u2013dissipation theorem ", INK, False),
          ("locks the noise to the kernel", OLIVE, True),
          (". Dissipation and fluctuation are two faces of one object \u2014 you "
           "cannot honestly keep one and drop the other.", INK, False)]])

    cards_slide(
        prs, menu, "Part 8 \u00b7 what K-theory silently deletes",
        [[("K-theory is the GLE with three terms mutilated", INK, True)]],
        [("\u2460  Memory \u2192 \u03b4(t)", "The kernel is collapsed to an instant: "
          "no history, no fast clock.", CLAY),
         ("\u2461  \u03bd\u209c constant in k", "Made spatially local \u2014 plain "
          "gradient diffusion, the elliptic field denied.", CLAY),
         ("\u2462  Noise \u2192 0", "FDT-linked fluctuation deleted: no backscatter, "
          "energy can only drain.", CLAY),
         ("What's left", "Exactly Smagorinsky: \u2202\u209c\u00fb = \u2119\u2096N\u0302 "
          "\u2212 \u03bd\u209c|k|\u00b2\u00fb. Not a bad approximation \u2014 a "
          "structural amputation.", MUTE)],
        accent=CLAY, cols=2)

    equation_slide(
        prs, menu, "Part 8 \u00b7 the closed model",
        [[("The repair: ", INK, True), ("projected-FDT closure", OLIVE, True)]],
        [[("\u2202\u209c\u00fb(k) = \u2119\u2096 N\u0302(\u00fb)(k) \u2212 "
           "\u03bd\u209c(k)|k|\u00b2 \u00fb(k) + f\u0302(k)", OLIVE, True)]],
        [[("(a) ", OLIVE, True), ("\u03bd\u209c(k)", INK, True),
          (" \u2014 nonlocal spectral eddy viscosity (Kraichnan cusp), "
           "allowed negative at low k \u2192 backscatter.", INK, False)],
         [("(b) ", OLIVE, True), ("f\u0302(k)", INK, True),
          (" \u2014 Leray-projected, FDT-linked noise: \u27e8f\u0302f\u0302\u27e9 = "
           "2\u03bd\u209c\u1d2e|k|\u00b2\u0398\u2119_ij \u2192 backscatter that "
           "respects div = 0.", INK, False)],
         [("(c) ", OLIVE, True), ("\u2119\u2096", INK, True),
          (" wraps everything \u2192 incompressible by construction. Limits: "
           "\u0394\u21920 \u2192 molecular NS; local/no-noise \u2192 Smagorinsky. A "
           "strict generalization.", INK, False)]])

    review_slide(prs, menu,
                 "Which term restores the fast clock?",
                 "Pair up: of the three things K-theory deletes, which single one "
                 "is the fast (elliptic) clock \u2014 and which restores backscatter?")

    # ===== SECTION 05 \u2014 THE BENCHMARK ===================================
    b = section_divider(
        prs, menu, "05",
        [[("The benchmark", CREAM, True)]],
        "A frozen-field a-priori test: 256\u00b2 DNS \u2192 sharp filter at k_c=32 "
        "\u2192 exact subgrid force \u2192 score three models on spectrum, "
        "divergence, and transfer.")
    anchors["05"] = b

    statement_slide(
        prs, menu, "The test that decides it",
        [[("Three models. ", INK, True), ("Three diagnostics.", TEAL, True)],
         [("Only one passes all three.", INK, True)]],
        sub="Smagorinsky (K-theory) vs. a spectrum-matched surrogate (the SPDE "
            "idea) vs. projected-FDT \u2014 graded against the exact subgrid force.")

    figure_slide(
        prs, menu, "Benchmark \u00b7 energy transfer T(k)",
        [[("Only projected-FDT reproduces ", INK, True),
          ("backscatter", OLIVE, True), (" \u2014 corr ", INK, True),
          ("1.000", OLIVE, True), (".", INK, True)]],
        f"{fig}/30_closure_transfer.png",
        stats=[("Smagorinsky over-dissipates", "cusp", CLAY),
               ("Surrogate transfer corr.", "0.907", OCHRE),
               ("Projected-FDT transfer corr.", "1.000", OLIVE)],
        caption="Truth shows dissipation (T<0) and backscatter (T>0); only "
                "projected-FDT reproduces both.")

    # The scorecard \u2014 clean table as soft cards, winning row highlighted
    slide = _blank(prs, CREAM)
    kicker(slide, "Benchmark \u00b7 the scorecard", color=OCHRE)
    tt = textbox(slide, MARGIN, 0.95, CW, 1.0)
    line_para(tt, [("Energy ", INK, True), ("and", OLIVE, True),
                   (" structure \u2014 that's the whole contribution.", INK, True)],
              size=27, font=HEAD, bold=True, first=True)
    # header
    cols_x = [MARGIN, MARGIN + 4.4, MARGIN + 7.0, MARGIN + 9.4]
    hdr = ["Model", "Spectrum E_m(k)", "RMS(\u2207\u00b7m)", "corr T(k)"]
    htf = textbox(slide, MARGIN, 2.25, CW, 0.4)
    rows = [
        ("Truth", "\u2014", "7.4\u00d710\u207b\u00b9\u2074", "1.000", MUTE, None),
        ("Smagorinsky (K-theory)", "over-dissipates", "1.8\u00d710\u207b\u00b9\u2074",
         "0.071", CLAY, "no backscatter"),
        ("Surrogate (SPDE)", "matches \u2713", "12.0", "0.907", CLAY,
         "divergence FAILS"),
        ("Projected-FDT", "matches \u2713", "1.7\u00d710\u207b\u00b9\u2074", "1.000",
         OLIVE, "all three pass"),
    ]
    top = 2.35
    rh = 0.92
    gap = 0.16
    # column headers
    for j, hx in enumerate(cols_x):
        cap = textbox(slide, hx + (0.3 if j == 0 else 0), top, 4.0, 0.35)
        line_para(cap, [(hdr[j], MUTE, True)], size=12.5, font=BODY, first=True)
    for i, (name, sp, dv, ct, acc, badge) in enumerate(rows):
        ry = top + 0.45 + i * (rh + gap)
        win = acc is OLIVE
        rrect(slide, MARGIN, ry, CW, rh, CARD if not win else RGBColor(0xEE, 0xF1, 0xE4),
              radius=0.10)
        rrect(slide, MARGIN, ry, 0.12, rh, acc if acc else MUTE, radius=0.10)
        nt = textbox(slide, cols_x[0] + 0.3, ry, 4.1, rh, MSO_ANCHOR.MIDDLE)
        line_para(nt, [(name, INK, True)], size=16, font=HEAD,
                  bold=True, first=True)
        if badge:
            line_para(nt, [("\u2192 " + badge, acc if acc else MUTE, False)],
                      size=11.5, font=BODY, space_before=2)
        for j, val in enumerate([sp, dv, ct], start=1):
            vt = textbox(slide, cols_x[j], ry, 2.4, rh, MSO_ANCHOR.MIDDLE)
            vcol = acc if (j == 3 and acc) else INK
            line_para(vt, [(val, vcol, j == 3)], size=16, font=HEAD, first=True)
    footer(slide, menu)
    notes(slide, "This is Fig 24 promoted to a closure benchmark. Walk the rows: "
                 "surrogate gets the energy but RMS divergence = 12 (catastrophic); "
                 "Smagorinsky is clean on divergence but corr 0.071 (no "
                 "backscatter); only projected-FDT passes all three. Energy AND "
                 "structure.")

    statement_slide(
        prs, menu, "What we proved",
        [[("K-theory isn't inaccurate.", INK, True)],
         [("It is ", INK, True), ("structurally blind", CLAY, True),
          (".", INK, True)]],
        sub="The blindness is precisely the Markov-delta collapse of the memory "
            "kernel plus deletion of the FDT noise. Restore both \u2014 with the "
            "Leray projector \u2014 and you get the exact repair, not a patch.")

    # ===== SECTION 06 \u2014 WHY IT MATTERS ==================================
    w = section_divider(
        prs, menu, "06",
        [[("Why it matters", CREAM, True)]],
        "One structural error, four real-world systems where the fast clock is the "
        "whole story \u2014 and K-theory quietly gets the physics backwards.")
    anchors["06"] = w

    cards_slide(
        prs, menu, "Application \u00b7 fast-clock physics",
        [[("Goldshtik\u2013Sorokin: a particle that ", INK, True),
          ("refuses to fall", TEAL, True)]],
        [("The phenomenon", "A heavy particle hangs suspended in a turbulent "
          "swirling flow. Small eddies self-organize a low-pressure vortex core; "
          "its gradient balances gravity.", TEAL),
         ("The requirement", "Exact elliptic pressure geometry + sustained swirl "
          "energy (backscatter) + div = 0 at all scales. Pure fast clock.", TEAL),
         ("K-theory", "\u03bd\u209c>0 smears the vortex core; no backscatter drains "
          "the swirl \u2192 the low-pressure pocket collapses. The particle falls.",
          CLAY),
         ("Projected-FDT", "\u2119 preserves the sharp core; negative \u03bd\u209c(k) "
          "feeds the swirl. The suspension is reproduced, not 'anomalous'.", OLIVE)],
        cols=2)

    cards_slide(
        prs, menu, "Application \u00b7 biomedical",
        [[("Liquid ventilation: ", INK, True),
          ("Dean vortices keep patients breathing", TEAL, True)]],
        [("The setup", "Dense perfluorocarbon is pumped through the bronchial "
          "tree. Oxygenation depends entirely on mixing.", TEAL),
         ("Fast clock", "At each bifurcation, flow separation spins up Dean "
          "vortices \u2014 fast, pressure-driven secondary currents that do the "
          "mixing.", TEAL),
         ("K-theory", "\u03bd\u209c>0 damps the vortices into smooth slug flow; "
          "spurious divergence spikes pressure \u2192 false 'suffocation + "
          "barotrauma' \u2014 both closure artifacts.", CLAY),
         ("Projected-FDT", "Backscatter sustains the vortices; \u2119 keeps the "
          "pressure field clean \u2192 correct: safe delivery, adequate "
          "oxygenation.", OLIVE)],
        cols=2)

    cards_slide(
        prs, menu, "Application \u00b7 geophysical",
        [[("Subglacial cavities: the fast clock ", INK, True),
          ("melts the ice", TEAL, True)]],
        [("The setup", "Pressurised meltwater flows over bedrock bumps under a "
          "glacier, carving turbulent cavities that melt ice from below.", TEAL),
         ("Fast clock", "Flow separates behind bumps into shedding vortices that "
          "trap heat and melt the roof; water-pressure spikes can lift the "
          "glacier.", TEAL),
         ("K-theory", "\u03bd\u209c>0 makes a smooth 'dead wake' \u2192 heat not "
          "trapped, melt under-predicted; spurious pressure \u2192 glacier looks "
          "safer than it is.", CLAY),
         ("Projected-FDT", "Live vortex shedding \u2192 realistic melt; exact mass "
          "conservation \u2192 correct effective pressure and sliding risk.", OLIVE)],
        cols=2)

    # ---- Part 9: the subglacial claim, now COMPUTED (not prose) -----------
    s = figure_slide(
        prs, menu, "Application \u00b7 geophysical (now computed)",
        [[("From a slide claim to a ", INK, True), ("computed", TEAL, True),
          (" benchmark \u2014 on a ", INK, True), ("real bed", TEAL, True)]],
        f"{fig}/35_subglacial_fields.png",
        caption="Part 9: 128\u00b2 penalized DNS of meltwater over a REAL BEDMAP1 "
                "Antarctic radar transect (BAS/SCAR, CC-BY-4.0; 220 km, 1757 m "
                "relief). Jet over the bumps, lee-side wakes, heat trapped against "
                "the warm bed. Only the geometry is measured \u2014 no instrument "
                "resolves subglacial DNS turbulence.")
    notes(s, "The reviewer's question answered: the subglacial 'application' used "
             "to be hand-written prose. It is now a run in the repo "
             "(run_subglacial.py --bed real), over a measured Antarctic bed.")

    s = figure_slide(
        prs, menu, "Application \u00b7 the test on real relief",
        [[("Same K-theory failure \u2014 ", INK, True),
          ("geometry-independent", CLAY, True)]],
        f"{fig}/36_subglacial_transfer.png",
        stats=[("Smagorinsky transfer corr", "\u22120.71", CLAY),
               ("projected-FDT transfer corr", "+0.45", OLIVE),
               ("surrogate \u2207\u00b7m (solenoidality)", "6.7", CLAY)],
        caption="Real BEDMAP1 bed: Smagorinsky stays dissipative at every scale "
                "(dead wake); the surrogate breaks \u2207\u00b7m = 0 (spurious bed "
                "pressure); only projected-FDT restores backscatter solenoidally. "
                "The idealized bed gives the same scorecard.")
    notes(s, "Swapping idealized sines for a measured asymmetric bed: a "
             "geometry-dependent artifact would move; a structural operator-level "
             "failure does not. It doesn't.")

    s = figure_slide(
        prs, menu, "Why it matters \u00b7 the mechanism, made literal",
        [[("One bump shapes the ", INK, True), ("global", TEAL, True),
          (" pressure,", INK, True)],
         [("only the ", INK, True), ("local", CLAY, True),
          (" temperature", INK, True)]],
        f"{fig}/39_elliptic_pressure.png",
        stats=[("pressure \u0394p beyond the bump", "79%", TEAL),
               ("temperature \u0394\u03b8 beyond the bump", "10%", CLAY)],
        caption="Same localized geometry source q applied to both operators: the "
                "elliptic pressure (\u2207\u207b\u00b2, like electrostatics) spans "
                "the cavity; the parabolic temperature (e^{t\u2207\u00b2}) stays a "
                "local blob. K-theory's \u03bd\u209c\u2207\u00b2 is structurally "
                "local \u2014 it cannot be the global pressure response.")
    notes(s, "The reviewer's own intuition made rigorous: solid geometry sets the "
             "boundary, the elliptic operator broadcasts it everywhere at once. "
             "Pressure is electrostatics; temperature is diffusion. One local "
             "diffusivity cannot be both \u2014 79% vs 10% far-field fraction.")

    # Otis \u2014 resolve the opening hook
    trigger_slide(
        prs, menu, "Application \u00b7 resolving the hook",
        [[("Back to Otis: the models couldn't", INK, True)],
         [("transfer energy ", INK, True), ("upscale", CLAY, True), (".", INK, True)]],
        "Rapid intensification is the fast clock dictating the slow one: subgrid "
        "convection organizes UPSCALE into the vortex. K-theory is bounded to "
        "T(k)\u22640 \u2014 it can only dissipate, so the ocean's heat was burned "
        "as 'subgrid friction' instead of spinning up the core.",
        "Callback to the opening clip \u2014 now the room has the mechanism.",
        accent=CLAY)

    cards_slide(
        prs, menu, "Application \u00b7 Hurricane Otis (2023)",
        [[("Projected-FDT lets a model ", INK, True),
          ("predict", OLIVE, True), (" rapid intensification", INK, True)]],
        [("Backscatter", "T(k)>0 at resolved scales \u2192 hot-tower bursts feed "
          "the mesoscale vortex \u2192 spin-up becomes possible.", OLIVE),
         ("Structure", "\u2119 keeps divergent outflow and inflow consistent \u2192 "
          "the extreme eye pressure drop is maintainable.", OLIVE),
         ("Scale-selective \u03bd\u209c(k)", "Doesn't kill the eyewall hot towers "
          "\u2192 latent-heat engine isn't artificially capped.", OLIVE),
         ("The Helmholtz view", "Models nailed the rotational steering, missed the "
          "divergent fast clock \u2014 the atmosphere's 'surrogate closure': right "
          "energy, wrong structure (Fig 29).", CLAY)],
        cols=2)

    # the universal pattern table
    slide = _blank(prs, CREAM)
    kicker(slide, "The universal pattern", color=TEAL)
    tt = textbox(slide, MARGIN, 0.95, CW, 1.0)
    line_para(tt, [("Four domains, ", INK, True), ("one", CLAY, True),
                   (" structural error, ", INK, True), ("one", OLIVE, True),
                   (" fix.", INK, True)], size=27, font=HEAD, bold=True, first=True)
    pat = [
        ("Goldshtik\u2013Sorokin", "vortex pressure core sustained by swirl",
         "core diffuses \u2192 particle falls"),
        ("PFC liquid ventilation", "Dean vortices at bronchial bifurcations",
         "vortices damped \u2192 mixing dies"),
        ("Subglacial cavity", "wake eddies melting ice behind bumps",
         "'dead wake' \u2192 melt under-predicted"),
        ("Hurricane Otis", "hot towers + upscale energy to vortex",
         "plumes smeared \u2192 RI not predicted"),
    ]
    top = 2.3
    rh = 0.82
    gap = 0.16
    for hx, htxt, hc in [(MARGIN + 0.32, "Phenomenon", MUTE),
                         (MARGIN + 3.9, "fast-clock structure", MUTE),
                         (MARGIN + 8.3, "K-theory failure", CLAY)]:
        ch_hdr = textbox(slide, hx, top - 0.04, 4.0, 0.32)
        line_para(ch_hdr, [(htxt, hc, True)], size=12, font=BODY, first=True)
    for i, (ph, st, fl) in enumerate(pat):
        ry = top + 0.35 + i * (rh + gap)
        rrect(slide, MARGIN, ry, CW, rh, CARD, radius=0.12)
        rrect(slide, MARGIN, ry, 0.12, rh, OCHRE, radius=0.12)
        a = textbox(slide, MARGIN + 0.32, ry, 3.5, rh, MSO_ANCHOR.MIDDLE)
        line_para(a, [(ph, INK, True)], size=15, font=HEAD, bold=True, first=True)
        bx = textbox(slide, MARGIN + 3.9, ry, 4.3, rh, MSO_ANCHOR.MIDDLE)
        line_para(bx, [(st, TEAL, False)], size=13.5, font=BODY, first=True)
        cx = textbox(slide, MARGIN + 8.3, ry, 3.2, rh, MSO_ANCHOR.MIDDLE)
        line_para(cx, [(fl, CLAY, False)], size=13.5, font=BODY, first=True)
    bt = textbox(slide, MARGIN, 6.55, CW, 0.5)
    line_para(bt, [("One fix:  ", INK, True),
                   ("\u2119 + FDT + \u03bd\u209c(k) \u2192 fast clock preserved "
                    "\u2192 physics captured.", OLIVE, True)],
              size=17, font=HEAD, first=True)
    footer(slide, menu)

    # ===== CLOSE ============================================================
    cards_slide(
        prs, menu, "Scope & honesty",
        [[("What we ", INK, True), ("do not", CLAY, True), (" claim", INK, True)]],
        [("Not universal", "No claim of a 'flawless' closure \u2014 turbulence "
          "closure is open.", CLAY),
         ("Not 3D regularity", "No Beale\u2013Kato\u2013Majda / wave-radiation-"
          "damping proof; those need rigorous analysis.", CLAY),
         ("Not operational (yet)", "No claim of superiority in a full weather/"
          "climate model.", CLAY),
         ("What IS claimed", "Exact structural constraints (\u2119, FDT, Galilean), "
          "removal of the isolated failure mode, two clean limits, and a "
          "falsifiable benchmark.", OLIVE)],
        accent=CLAY, cols=2)

    statement_slide(
        prs, menu, "The takeaway",
        [[("The error was never ", INK, True), ("missing data", MUTE, True),
          (".", INK, True)],
         [("It was destroying the fast clock.", INK, True)]],
        sub="Restore the memory kernel (including backscatter) and the FDT-linked, "
            "Leray-projected noise, and K-theory's structural blindness is cured "
            "\u2014 exactly, and reproducibly.")

    poll_slide(prs, menu,
               "One word: what will you remember from this?",
               ["Two clocks", "Structural, not data", "Backscatter",
                "Projection / div = 0"],
               code="TWOCLOCKS", tool="Mentimeter")

    # thank-you / repro
    slide = _blank(prs, CREAM)
    circle(slide, -1.6, 4.6, 4.4, SAND)
    circle(slide, 10.8, -1.8, 4.6, RGBColor(0xEE, 0xE7, 0xD6))
    rrect(slide, MARGIN, 2.0, 0.9, 0.16, TEAL, radius=0.5)
    tf = textbox(slide, MARGIN, 2.35, 11.2, 2.0)
    line_para(tf, [("Thank you.", TEAL, True)], size=58, font=HEAD, bold=True,
              first=True)
    sub = textbox(slide, MARGIN, 4.0, 11.0, 2.2)
    line_para(sub, [("Every figure and number here is regenerated from source.",
                     INK, False)], size=20, font=BODY, first=True)
    line_para(sub, [("python run_closure.py", TEAL, True),
                    ("   \u2192   the benchmark that decides it.", MUTE, False)],
              size=20, font=MONO, space_before=10)
    line_para(sub, [("Abdelrahman Saifelden", INK, True),
                    ("   \u00b7   Mechatronics Engineering, AIET Alexandria",
                     MUTE, False)], size=18, font=BODY, space_before=14)
    footer(slide, menu)
    notes(slide, "Close on reproducibility. Offer to run the benchmark live if "
                 "anyone doubts a number \u2014 it's one command.")

    # ---- wire the hyperlinked menu ----------------------------------------
    order = ["01", "02", "03", "04", "05", "06"]
    for card, key in zip(menu_cards, order):
        if key in anchors:
            card.click_action.target_slide = anchors[key]

    out_path = "Two_Clocks_Presentation.pptx"
    prs.save(out_path)
    print(f"Presentation saved: {out_path}")
    print(f"  {len(prs.slides)} slides")


if __name__ == "__main__":
    main()
