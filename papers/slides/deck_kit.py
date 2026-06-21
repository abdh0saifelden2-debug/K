#!/usr/bin/env python3
"""Generate the talk decks (v2) for Papers 1-4, following researched craft:

  * Assertion-Evidence (Alley, Penn State): every content slide is a ONE-SENTENCE
    assertion headline supported by VISUAL evidence (a real figure or a built
    result/test panel) — no bullet lists; secondary detail lives in the speaker
    notes.
  * ABT narrative (Olson): each deck is an And-But-Therefore story.
  * Write-to-speak SCRIPT in the notes pane of every slide (mirrored to SCRIPT.md),
    carrying the explanation + one analogy per hard idea.
  * 60-30-10 colour, WCAG-safe: soft warm off-white background, one per-paper
    identity hue, an accent reserved for key results; status always carries a text
    label (never colour alone).

Outputs per paper:  paperN_*.pptx  +  paperN_SCRIPT.md
Needs python-pptx (+ Pillow for image aspect ratios).
"""
from __future__ import annotations
import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

try:
    from PIL import Image
    def aspect(p):
        with Image.open(p) as im:
            return im.width / im.height
except Exception:                                  # pragma: no cover
    def aspect(p):
        return 1.6

REPO = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
GL = os.path.join(REPO, "glaciers")
OUT = os.path.dirname(os.path.abspath(__file__))

W, H = Inches(13.333), Inches(7.5)
BG = RGBColor(0xF6, 0xF7, 0xF9)        # soft warm off-white (60%)
INK = RGBColor(0x1E, 0x23, 0x2B)       # softened near-black
SUB = RGBColor(0x55, 0x5C, 0x66)       # secondary text
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
VER = RGBColor(0x1B, 0x7F, 0x37)
FAL = RGBColor(0xB3, 0x26, 0x1A)
OPN = RGBColor(0xB5, 0x7A, 0x00)


def tint(c, f):
    """Lighten colour c toward white by fraction f (0..1)."""
    return RGBColor(int(c[0] + (255 - c[0]) * f), int(c[1] + (255 - c[1]) * f),
                    int(c[2] + (255 - c[2]) * f))


def _run(p, text, size, color, bold=False, italic=False, font="Calibri"):
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.name = font; r.font.color.rgb = color
    return r


def _box(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(l, t, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    return tf


def _rect(slide, l, t, w, h, color, shape=MSO_SHAPE.RECTANGLE, line=None):
    sp = slide.shapes.add_shape(shape, l, t, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


class Deck:
    def __init__(self, accent, short, title_short):
        self.prs = Presentation(); self.prs.slide_width = W; self.prs.slide_height = H
        self.accent = accent; self.short = short; self.title_short = title_short
        self.n = 0; self.blank = self.prs.slide_layouts[6]
        self.script = []                            # (slide_title, notes) for SCRIPT.md

    def _new(self):
        s = self.prs.slides.add_slide(self.blank)
        s.background.fill.solid(); s.background.fill.fore_color.rgb = BG
        return s

    def _notes(self, slide, head, text):
        slide.notes_slide.notes_text_frame.text = text
        self.script.append((head, text))

    def _footer(self, slide):
        self.n += 1
        tf = _box(slide, Inches(0.5), Inches(7.04), Inches(10.0), Inches(0.36))
        _run(tf.paragraphs[0], f"{self.short}   ·   github.com/abdh0saifelden2-debug/K", 9.5, SUB)
        tf2 = _box(slide, Inches(12.3), Inches(7.04), Inches(0.7), Inches(0.36))
        tf2.paragraphs[0].alignment = PP_ALIGN.RIGHT
        _run(tf2.paragraphs[0], str(self.n), 9.5, SUB)

    def _headline(self, slide, text):
        _rect(slide, 0, 0, W, Inches(1.5), tint(self.accent, 0.86))
        _rect(slide, 0, Inches(1.5), W, Pt(2.5), self.accent)
        tf = _box(slide, Inches(0.55), Inches(0.12), Inches(12.2), Inches(1.25), MSO_ANCHOR.MIDDLE)
        _run(tf.paragraphs[0], text, 25, INK, bold=True)

    def _cards(self, slide, cards, top):
        """Result/test 'evidence' panel: metric cards (big number + label)."""
        n = len(cards); gap = Inches(0.25)
        total = Inches(12.2); cw = (total - gap * (n - 1)) / n
        for i, (big, label) in enumerate(cards):
            left = Inches(0.55) + i * (cw + gap)
            _rect(slide, left, top, cw, Inches(1.5), tint(self.accent, 0.93),
                  shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            tf = _box(slide, left + Inches(0.12), top + Inches(0.12), cw - Inches(0.24),
                      Inches(1.26), MSO_ANCHOR.MIDDLE)
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            _run(tf.paragraphs[0], big, 23, self.accent, bold=True)
            p = tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER
            _run(p, label, 12, SUB)

    def _image(self, slide, path, left, top, maxw, maxh):
        ar = aspect(path); w = maxw; h = Emu(int(w / ar))
        if h > maxh:
            h = maxh; w = Emu(int(h * ar))
        self._cur = slide.shapes.add_picture(path, left + (maxw - w) / 2,
                                             top + (maxh - h) / 2, width=w, height=h)

    def _status(self, slide, status, top):
        color = {"VERIFIED": VER, "FALSIFIED": FAL, "OPEN": OPN}[status[0]]
        _rect(slide, Inches(0.55), top, Inches(0.22), Inches(0.34), color)
        tf = _box(slide, Inches(0.85), top - Inches(0.03), Inches(11.8), Inches(0.42), MSO_ANCHOR.MIDDLE)
        _run(tf.paragraphs[0], f"{status[0]}  —  {status[1]}", 12.5, color, bold=True)

    # ---- slide kinds ---- #
    def title(self, title, thesis, author, script):
        s = self._new()
        _rect(s, 0, 0, W, Inches(2.7), self.accent)
        _rect(s, 0, Inches(2.7), W, Pt(4), tint(self.accent, 0.4))
        tf = _box(s, Inches(0.7), Inches(0.55), Inches(12.0), Inches(2.0), MSO_ANCHOR.MIDDLE)
        _run(tf.paragraphs[0], title, 30, WHITE, bold=True)
        tf2 = _box(s, Inches(0.7), Inches(3.05), Inches(12.0), Inches(2.6))
        _run(tf2.paragraphs[0], "The story in one line", 13, self.accent, bold=True)
        p = tf2.add_paragraph(); p.space_before = Pt(6); _run(p, thesis, 19, INK)
        tf3 = _box(s, Inches(0.7), Inches(6.2), Inches(12.0), Inches(1.0))
        _run(tf3.paragraphs[0], author, 13, SUB)
        p = tf3.add_paragraph(); _run(p, "code & data: github.com/abdh0saifelden2-debug/K", 12, SUB)
        self.n += 1; self._notes(s, title, script)

    def analogy(self, head, lines, script, big=None):
        s = self._new(); self._headline(s, head)
        tf = _box(s, Inches(0.7), Inches(2.1), Inches(11.9), Inches(3.6))
        for i, ln in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(14)
            _run(p, ln, 21 if i == 0 else 18, INK if i == 0 else SUB, italic=(i == 0))
        if big:
            bx = _box(s, Inches(0.7), Inches(5.4), Inches(11.9), Inches(1.4), MSO_ANCHOR.MIDDLE)
            _run(bx.paragraphs[0], big, 16, self.accent, bold=True)
        self._footer(s); self._notes(s, head, script)

    def assert_slide(self, head, script, image=None, caption=None, cards=None,
                     support=None, status=None):
        s = self._new(); self._headline(s, head)
        has_img = image and os.path.exists(image)
        y = Inches(1.75)
        if support:
            tf = _box(s, Inches(0.55), y, Inches(12.2), Inches(0.55))
            _run(tf.paragraphs[0], support, 16, SUB); y = Inches(2.35)
        if has_img and cards:
            self._image(s, image, Inches(0.55), y, Inches(7.4), Inches(4.1))
            # cards stacked on the right
            top = y; ch = Inches(0.95); gap = Inches(0.18)
            for i, (big, label) in enumerate(cards[:4]):
                t = top + i * (ch + gap)
                _rect(s, Inches(8.2), t, Inches(4.55), ch, tint(self.accent, 0.93),
                      shape=MSO_SHAPE.ROUNDED_RECTANGLE)
                cf = _box(s, Inches(8.35), t + Inches(0.06), Inches(4.25), ch - Inches(0.12), MSO_ANCHOR.MIDDLE)
                _run(cf.paragraphs[0], big, 19, self.accent, bold=True)
                p = cf.add_paragraph(); _run(p, label, 12, SUB)
            if caption:
                cf = _box(s, Inches(0.55), y + Inches(4.15), Inches(7.4), Inches(0.4))
                _run(cf.paragraphs[0], caption, 11, SUB, italic=True)
        elif has_img:
            self._image(s, image, Inches(2.4), y, Inches(8.5), Inches(4.4))
            if caption:
                cf = _box(s, Inches(0.55), Inches(6.35), Inches(12.2), Inches(0.4))
                cf.paragraphs[0].alignment = PP_ALIGN.CENTER
                _run(cf.paragraphs[0], caption, 11, SUB, italic=True)
        elif cards:
            self._cards(s, cards, Inches(2.9))
            if caption:
                cf = _box(s, Inches(0.55), Inches(4.7), Inches(12.2), Inches(0.5))
                cf.paragraphs[0].alignment = PP_ALIGN.CENTER
                _run(cf.paragraphs[0], caption, 12, SUB, italic=True)
        if status:
            self._status(s, status, Inches(6.5))
        self._footer(s); self._notes(s, head, script)

    def scope(self, verified, falsified, openitems, script):
        s = self._new(); self._headline(s, "What is verified, what is falsified, what is still open")
        for j, (head, color, items) in enumerate([("VERIFIED", VER, verified),
                                                   ("FALSIFIED", FAL, falsified),
                                                   ("OPEN", OPN, openitems)]):
            left = Inches(0.55 + j * 4.13)
            _rect(s, left, Inches(1.9), Inches(3.9), Inches(0.5), color)
            hf = _box(s, left + Inches(0.15), Inches(1.93), Inches(3.6), Inches(0.44), MSO_ANCHOR.MIDDLE)
            _run(hf.paragraphs[0], head, 14, WHITE, bold=True)
            _rect(s, left, Inches(2.4), Inches(3.9), Inches(4.2), tint(color, 0.92))
            bf = _box(s, left + Inches(0.18), Inches(2.55), Inches(3.55), Inches(3.9))
            for i, it in enumerate(items):
                p = bf.paragraphs[0] if i == 0 else bf.add_paragraph()
                p.space_after = Pt(7); _run(p, "•  " + it, 12.5, INK)
        self._footer(s); self._notes(s, "Scope", script)

    def save(self, fname):
        path = os.path.join(OUT, fname); self.prs.save(path)
        # write SCRIPT.md
        md = os.path.join(OUT, fname.replace(".pptx", "_SCRIPT.md"))
        with open(md, "w") as f:
            f.write(f"# Speaker script — {self.title_short}\n\n")
            f.write("Write-to-speak narration, one block per slide (also embedded in the "
                    "deck's notes pane). Say it in your own words; let the slide carry the visual.\n\n")
            for i, (head, note) in enumerate(self.script, 1):
                f.write(f"## Slide {i} — {head}\n\n{note}\n\n")
        return path, md


def fig(*p):
    return os.path.join(GL, *p)
